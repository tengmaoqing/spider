from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import os
import re

def set_text_frame_font(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
        #   print(run)
          run.text = run.text.replace('柚墨', '清静')
          run.text = run.text.replace('柚小墨', '小清静')
          run.text = re.sub(r'Yomoer', 'baidu1', run.text, flags=re.IGNORECASE)
          run.text = run.text.replace('YOZOPPT', '清静')

def check_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shape_in_group in shape.shapes:
            check_shape(shape_in_group)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for cell in shape.table.iter_cells():
            text_frame = cell.text_frame
            set_text_frame_font(text_frame)
    else:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            set_text_frame_font(text_frame)

def start(file, dist):
    print(f'Processing file: {file}')
    try:
        prs = Presentation(file)
        for index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                check_shape(shape)
        # target = f'temp/{file}'
        # dir = os.path.dirname(target)
        # if not os.path.exists(dir):
        #     os.makedirs(dir)
        prs.save(dist)
    except Exception:
        print(f'error {file}')
    finally:
        return dist

# for file in glob.glob('download/**/*.pptx', recursive=True):
#     print(f'Processing file: {file}')
#     try:
#         prs = Presentation(file)
#         for index, slide in enumerate(prs.slides):
#             for shape in slide.shapes:
#                 check_shape(shape)
#         target = f'temp/{file}'
#         dir = os.path.dirname(target)
#         if not os.path.exists(dir):
#             os.makedirs(dir)
#         prs.save(target)
#     except Exception:
#         print(f'error {file}')